import sys
from pathlib import Path
import argparse
import re
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches

# Ensure we can import the LLM util
sys.path.insert(0, str(Path(__file__).resolve().parents[2]))
from app_dms_global.utils.llm import ask

def parse_slides_from_markdown(md: str):
    """
    Only H1 (# ) and H2 (## ) become slides.
    Everything else stays in the body.
    """
    slides = []
    current_title = None
    current_body = []

    for line in md.splitlines():
        if line.startswith("## "):
            if current_title is not None:
                slides.append((current_title, current_body))
            current_title = line[3:].strip()
            current_body = []
        elif line.startswith("# "):
            if current_title is None:
                current_title = line[2:].strip()
                current_body = []
            else:
                current_body.append(line)
        else:
            current_body.append(line)

    if current_title is not None:
        slides.append((current_title, current_body))

    return slides

def clean_slides(slides):
    """
    Drop slides whose title is internal/technical.
    """
    out = []
    for title, body in slides:
        t = title.lower().strip()
        if t in ("insight","key insights","document insights","insights from uploaded documents"):
            continue
        if t.startswith("from ") or t.endswith((".md",".pdf",".docx")):
            continue
        text = "\n".join(body).strip()
        if not text and out:
            continue
        out.append((title, body))
    return out

def load_insight_bullets(insights_path: Path):
    """
    Return up to 3 bullets from key_insights.md
    """
    if not insights_path.exists():
        return []
    bullets = []
    for line in insights_path.read_text(encoding="utf-8").splitlines():
        l = line.strip()
        if l.startswith(("- ", "• ")):
            bullets.append(l.lstrip("-• ").strip())
            if len(bullets) >= 3:
                break
    return bullets

def parse_insight_sources(insights_path: Path):
    """
    Return sorted list of unique filenames from 'Source: xxx.md' lines.
    """
    if not insights_path.exists():
        return []
    text = insights_path.read_text(encoding="utf-8")
    sources = re.findall(r"Source:\s*([^\n]+\.md)", text)
    return sorted(set(sources))

def deck_md_to_ppt(project_name, audience="investors", sector=None, territory=None):
    BASE = Path(f"projects/{project_name}")
    MD   = BASE / "deck/deck.md"
    INS  = BASE / "insights/key_insights.md"
    CHART= BASE / "financial/charts/revenue.png"
    OUT  = Path(f"outputs/{project_name}-{audience}.pptx")
    OUT.parent.mkdir(parents=True, exist_ok=True)

    if not MD.exists():
        raise FileNotFoundError(f"{MD} not found")

    md = MD.read_text(encoding="utf-8")
    slides = clean_slides(parse_slides_from_markdown(md))
    print(f"📊 Will create {len(slides)} slides")

    tpl = Path("templates/Brand_Template.pptx")
    prs = Presentation(str(tpl)) if tpl.exists() else Presentation()
    layouts = prs.slide_layouts
    L_TITLE   = layouts[0]
    L_CONTENT = layouts[3] if len(layouts)>3 else layouts[0]
    L_FIN     = layouts[4] if len(layouts)>4 else L_CONTENT

    today = datetime.now().strftime("%Y-%m-%d")

    for idx, (title, body) in enumerate(slides):
        # choose layout
        if idx == 0:
            layout = L_TITLE
        elif any(k in title.lower() for k in ("financial","revenue","ebitda")):
            layout = L_FIN
        else:
            layout = L_CONTENT

        slide = prs.slides.add_slide(layout)

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = title

        # Subtitle on cover slide: territory
        if idx == 0 and territory:
            try:
                ph2 = slide.placeholders[1]
                ph2.text = territory
            except Exception:
                tb = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.5))
                tb.text_frame.text = territory

        # Project Overview
        if title.lower() == "project overview":
            # body text if available
            sentence = next((l.strip() for l in body if l.strip() and not l.startswith("#")), None)
            if not sentence and sector and territory:
                prompt = f"Write a one-sentence project overview for a {sector} project in {territory}."
                sentence = ask(prompt)
            if sentence:
                for ph in slide.placeholders:
                    if ph.has_text_frame and ph != slide.shapes.title:
                        ph.text = sentence
                        break
            print("  ✓ Project Overview set")
            continue

        # Executive Summary
        if title.lower() == "executive summary":
            exec_bullets = []
            if sector:    exec_bullets.append(f"Sector: {sector}")
            if territory: exec_bullets.append(f"Territory: {territory}")
            exec_bullets.append(f"Date: {today}")
            exec_bullets += load_insight_bullets(INS)
            for ph in slide.placeholders:
                if ph.has_text_frame and ph != slide.shapes.title:
                    ph.text = "\n".join(f"• {b}" for b in exec_bullets)
                    break
            print("  ✓ Executive Summary set")
            continue

        # Normal slides
        bullets = [l.lstrip("-• ").strip() for l in body if l.strip().startswith(("- ","• "))]
        if bullets:
            for ph in slide.placeholders:
                if ph.has_text_frame and ph != slide.shapes.title:
                    ph.text = "\n".join(f"• {b}" for b in bullets[:8])
                    break

        # Financial chart
        if layout == L_FIN and CHART.exists():
            slide.shapes.add_picture(str(CHART), Inches(1), Inches(2), width=Inches(8))

        print(f"  ✓ Slide {idx+1}: {title}")

    # Appendix — Sources
    sources = parse_insight_sources(INS)
    if sources:
        slide = prs.slides.add_slide(L_CONTENT)
        if slide.shapes.title:
            slide.shapes.title.text = "Appendix — Sources"
        for ph in slide.placeholders:
            if ph.has_text_frame and ph != slide.shapes.title:
                ph.text = "\n".join(f"• {s}" for s in sources)
                break
        print(f"  ✓ Appendix slide added with {len(sources)} sources")

    prs.save(OUT)
    print(f"✅ Saved PPT: {OUT} ({len(prs.slides)} slides)")
    return str(OUT)

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Convert deck.md to PPTX")
    parser.add_argument("--project", required=True)
    parser.add_argument("--audience", default="investors")
    parser.add_argument("--sector", required=True)
    parser.add_argument("--territory", required=True)
    args = parser.parse_args()

    deck_md_to_ppt(args.project, args.audience, args.sector, args.territory)