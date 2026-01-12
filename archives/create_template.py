# create_template.py
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

def create_basic_template():
    """Create a basic PowerPoint template if none exists."""
    
    template_path = Path("templates/Brand_Template.pptx")
    template_path.parent.mkdir(parents=True, exist_ok=True)
    
    # Create a new presentation
    prs = Presentation()
    
    # Slide layouts
    layouts = prs.slide_layouts
    
    # 1. Title Slide (layout 0)
    title_slide = prs.slides.add_slide(layouts[0])
    title_slide.shapes.title.text = "Title Slide"
    title_slide.placeholders[1].text = "Subtitle"
    
    # 2. Section Header (layout 1)
    section_slide = prs.slides.add_slide(layouts[1])
    section_slide.shapes.title.text = "Section Header"
    
    # 3. Title and Content (layout 3)
    content_slide = prs.slides.add_slide(layouts[3])
    content_slide.shapes.title.text = "Content Title"
    content_slide.placeholders[1].text = "• Bullet point 1\n• Bullet point 2\n• Bullet point 3"
    
    # 4. Two Content (layout 4) - using as financial slide
    financial_slide = prs.slides.add_slide(layouts[4])
    financial_slide.shapes.title.text = "Financial Overview"
    financial_slide.placeholders[1].text = "Revenue Projection\n• Year 1: $1.25M\n• Year 2: $1.68M\n• Year 3: $2.25M"
    
    # Save the template
    prs.save(template_path)
    
    print(f"✅ Created basic template: {template_path}")
    return template_path

if __name__ == "__main__":
    create_basic_template()